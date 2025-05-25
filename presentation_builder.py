"""
Presentation Builder Module - LightYearAI

This module provides functionality to create PowerPoint presentations from text content or document files.
It handles extraction of text from different file types, content analysis, and PPTX generation.
"""

import os
import uuid
import tempfile
from pathlib import Path
import re
from io import BytesIO

# Document parsing libraries
import docx
import pdfplumber
import fitz  # PyMuPDF

# Presentation creation library
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE

# Constants for presentation generation
MAX_SLIDES = 15
MAX_BULLET_POINTS_PER_SLIDE = 5
MIN_WORDS_FOR_BULLET = 3
MAX_WORDS_FOR_TITLE = 8


def extract_text_from_file(file_path, file_type=None):
    """
    Extract text content from a file based on its type.
    
    Args:
        file_path (str): Path to the file
        file_type (str, optional): File type override. If None, determined from file extension.
    
    Returns:
        str: Extracted text content
    """
    if file_type is None:
        # Determine file type from extension
        file_ext = Path(file_path).suffix.lower()
        if file_ext == '.txt':
            file_type = 'txt'
        elif file_ext == '.docx':
            file_type = 'docx'
        elif file_ext == '.pdf':
            file_type = 'pdf'
        else:
            raise ValueError(f"Unsupported file type: {file_ext}")
    
    # Extract text based on file type
    if file_type == 'txt':
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        except UnicodeDecodeError:
            # Try alternative encodings
            try:
                with open(file_path, 'r', encoding='latin-1') as f:
                    return f.read()
            except Exception as e:
                print(f"Error reading txt file: {str(e)}")
                return "Error extracting text from file."
    
    elif file_type == 'docx':
        try:
            doc = docx.Document(file_path)
            return '\n'.join([paragraph.text for paragraph in doc.paragraphs if paragraph and paragraph.text])
        except Exception as e:
            print(f"Error extracting text from docx: {str(e)}")
            return "Error extracting text from DOCX file."
    
    elif file_type == 'pdf':
        # Try with pdfplumber first
        extracted_text = ""
        try:
            text = []
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    if page is None:
                        continue
                    try:
                        page_text = page.extract_text()
                        if page_text:
                            text.append(page_text)
                    except Exception as page_err:
                        print(f"Error extracting text from PDF page: {str(page_err)}")
                        continue
            
            extracted_text = '\n'.join(text)
            
            # If pdfplumber extraction is empty or limited, try PyMuPDF as fallback
            if not extracted_text or len(extracted_text.strip()) < 50:
                raise Exception("Limited text extracted, trying fallback method")
                
            return extracted_text
                
        except Exception as e:
            print(f"pdfplumber extraction issue: {str(e)}, trying PyMuPDF")
            
            # PyMuPDF as fallback
            try:
                text = []
                doc = fitz.open(file_path)
                for page_num in range(len(doc)):
                    try:
                        page = doc.load_page(page_num)
                        if page:
                            page_text = page.get_text()
                            if page_text:
                                text.append(page_text)
                    except Exception as page_err:
                        print(f"Error with PyMuPDF page {page_num}: {str(page_err)}")
                        continue
                doc.close()
                
                extracted_text = '\n'.join(text)
                if not extracted_text or len(extracted_text.strip()) < 10:
                    return "Could not extract meaningful text from the PDF file."
                return extracted_text
            except Exception as mupdf_err:
                print(f"PyMuPDF extraction failed: {str(mupdf_err)}")
                return "Error extracting text from PDF file."
    
    else:
        raise ValueError(f"Unsupported file type: {file_type}")


def analyze_content(text, max_slides=MAX_SLIDES):
    """
    Analyze and structure text content for presentation slides.
    This function identifies sections and key points to organize into slides.
    
    Args:
        text (str): The text content to analyze
        max_slides (int): Maximum number of slides to generate
    
    Returns:
        dict: Structured content for presentation slides
    """
    # Safety check for None or invalid input
    if text is None:
        print("Warning: analyze_content received None as text input")
        text = "No content was provided for analysis."
    
    # Ensure text is a string
    if not isinstance(text, str):
        try:
            text = str(text)
            print(f"Converted non-string input to string in analyze_content, type was: {type(text)}")
        except:
            print(f"Could not convert input to string, type: {type(text)}")
            text = "Error converting content to proper format for analysis."
    
    # Initial cleaning of the text
    text = text.strip()
      # Check if the text is too short
    if len(text) < 100:
        return {
            'title': 'Simple Presentation',
            'slides': [
                {
                    'title': 'Key Points',
                    'points': [text] if text else ['No content provided']
                },
                {
                    'title': 'Summary',
                    'points': ['This is a brief summary of the key points.', 
                              'Consider adding more content for a more detailed presentation.',
                              'Contact LightYearAI support if you need assistance.']
                }
            ]
        }
    
    # Try to identify sections based on common patterns
    sections = []
    
    try:
        # Pattern 1: Numbered sections like "1. Section Title" or "Section 1: Title"
        section_pattern1 = re.compile(r'(?:\n|^)(?:\d+\.\s+|\bSection\s+\d+[:\.\s]+)([^\n]+)(?:\n|$)', re.IGNORECASE)
        matches = section_pattern1.findall(text)
        if len(matches) > 2:  # Only use this pattern if we find enough matches
            sections = matches[:max_slides-1]  # Reserve one slide for intro
        
        # Pattern 2: Headers with newlines or all caps
        if not sections:
            section_pattern2 = re.compile(r'(?:\n\s*\n|\n)([A-Z][^.\n]{3,30}[:.!?]?)(?:\n|$)')
            matches = section_pattern2.findall(text)
            if len(matches) > 2:
                sections = matches[:max_slides-1]
        
        # Pattern 3: Split by paragraphs if sections aren't found
        if not sections:
            paragraphs = [p.strip() for p in re.split(r'\n\s*\n', text) if p.strip()]
            if len(paragraphs) <= max_slides:
                # Each paragraph becomes a slide
                return {
                    'title': extract_title(text),
                    'slides': [
                        {
                            'title': create_slide_title(p),
                            'points': extract_bullet_points(p)
                        } for p in paragraphs
                    ]
                }
            else:
                # Combine paragraphs if there are too many
                grouped_paragraphs = []
                current_group = []
                total_length = 0
                
                for p in paragraphs:
                    if total_length > 500:  # Approx chars per slide
                        grouped_paragraphs.append('\n\n'.join(current_group))
                        current_group = [p]
                        total_length = len(p)
                    else:
                        current_group.append(p)
                        total_length += len(p)
                
                if current_group:
                    grouped_paragraphs.append('\n\n'.join(current_group))
                
                return {
                    'title': extract_title(text),
                    'slides': [
                        {
                            'title': create_slide_title(p),
                            'points': extract_bullet_points(p)
                        } for p in grouped_paragraphs[:max_slides]
                    ]
                }
        
        # If sections were found, extract content for each section
        slide_content = []
        
        # Get content between sections
        for i, section in enumerate(sections):
            section_title = section.strip()
            
            # Find the start position of this section
            start_pos = text.find(section_title)
            if start_pos == -1:  # Try again with more flexible search
                for s in [section_title, section_title.upper(), section_title.lower()]:
                    start_pos = text.find(s)
                    if start_pos != -1:
                        break
            
            if start_pos != -1:
                # Find the end of this section (start of next section or end of text)
                end_pos = len(text)
                if i < len(sections) - 1:
                    next_section = sections[i + 1].strip()
                    next_pos = text.find(next_section, start_pos + len(section_title))
                    if next_pos != -1:
                        end_pos = next_pos
                
                # Extract section content
                section_content = text[start_pos + len(section_title):end_pos].strip()
                
                # Create slide data
                slide_content.append({
                    'title': section_title,
                    'points': extract_bullet_points(section_content)
                })
        
        # If we couldn't extract content properly, fallback to a simple approach
        if not slide_content:
            paragraphs = text.split('\n\n')
            slide_content = [
                {
                    'title': create_slide_title(p),
                    'points': extract_bullet_points(p)
                } for p in paragraphs[:max_slides]
            ]
        
        # Prepare the presentation structure
        presentation_data = {
            'title': extract_title(text),
            'slides': slide_content[:max_slides]
        }
        
        return presentation_data
    except Exception as e:
        print(f"Error in analyze_content: {str(e)}")
        # Return a simple fallback presentation structure
        return {
            'title': 'Presentation',
            'slides': [
                {
                    'title': 'Main Content',
                    'points': ['There was an error analyzing the content.', 
                              'A simple presentation has been created instead.', 
                              'Please try again with different content.']
                }
            ]
        }


def extract_title(text):
    """Extract a title from the text content"""
    # Get the first non-empty line as potential title
    lines = text.split('\n')
    first_line = next((line for line in lines if line.strip()), 'Presentation')
    
    # Clean up the title
    title = first_line.strip().rstrip('.,:;')
    
    # Limit title length
    words = title.split()
    if len(words) > MAX_WORDS_FOR_TITLE:
        title = ' '.join(words[:MAX_WORDS_FOR_TITLE]) + '...'
    
    return title


def create_slide_title(text):
    """Create a slide title from a paragraph"""
    # Get the first line or first sentence
    if '\n' in text:
        first_line = text.split('\n')[0]
    else:
        sentences = re.split(r'(?<=[.!?])\s+', text)
        first_line = sentences[0] if sentences else text
    
    # Clean and limit the title
    title = first_line.strip().rstrip('.,:;')
    
    words = title.split()
    if len(words) > MAX_WORDS_FOR_TITLE:
        title = ' '.join(words[:MAX_WORDS_FOR_TITLE]) + '...'
    
    return title


def extract_bullet_points(text, max_points=MAX_BULLET_POINTS_PER_SLIDE):
    """Extract bullet points from text"""
    # Check if the text already has bullet points
    bullet_pattern = re.compile(r'(?:^|\n)[•\-\*\+]\s+(.+)(?:\n|$)')
    existing_bullets = bullet_pattern.findall(text)
    
    if existing_bullets:
        return existing_bullets[:max_points]
    
    # Extract sentences as bullet points
    sentences = re.split(r'(?<=[.!?])\s+', text.strip())
    
    # Filter out very short sentences
    bullets = []
    for sentence in sentences:
        sentence = sentence.strip()
        if sentence and len(sentence.split()) >= MIN_WORDS_FOR_BULLET:
            bullets.append(sentence)
    
    # Limit the number of bullet points
    return bullets[:max_points]


def enhance_slide_content_with_ai(slides, model=None):
    """
    Enhance slide content with AI-generated content when available
    Replaces placeholder-like content with more substantial information
    
    Args:
        slides (list): List of slide data dictionaries
        model: Optional GenAI model to use for enhancement
        
    Returns:
        list: Enhanced slides list
    """
    if not model:
        return slides
        
    try:
        # Ensure we have at least 5 slides for a good presentation
        if len(slides) < 5:
            # Generate additional slides with generic topics
            generic_topics = ["Background Information", "Key Benefits", "Implementation Strategy", 
                             "Future Opportunities", "Case Studies", "Market Analysis", "Recommendations"]
            
            for i in range(len(slides), 5):
                topic = generic_topics[i % len(generic_topics)]
                slides.append({
                    'title': topic,
                    'bullets': generate_bullet_points_from_title(topic)
                })
            
        for i, slide in enumerate(slides):
            title = slide.get('title', '')
            points = slide.get('bullets', [])
            
            # Skip if we don't have placeholder-like content and have sufficient bullet points
            placeholder_indicators = ["discuss the", "outline the", "describe the", "explain the", "list the"]
            has_placeholders = False
            
            # Check if we have enough high-quality content
            has_sufficient_content = len(points) >= 3
            
            for point in points:
                lower_point = point.lower()
                if any(indicator in lower_point for indicator in placeholder_indicators):
                    has_placeholders = True
                    has_sufficient_content = False
                    break
                    
            if not has_placeholders and has_sufficient_content:
                continue
                
            # If we found placeholder content, generate real content with AI
            prompt = f"""
            Generate 4 specific, factual bullet points about the topic: "{title}"
            
            Rules:
            - Each bullet point should be a specific fact, statistic, or concrete piece of information
            - Do NOT use placeholder language like "Discuss..." or "Outline..."
            - Keep each bullet point under 15 words if possible
            - Start each point with a strong action verb or specific data point
            - Return only the bullet points, one per line, with no numbering or symbols
            
            Example (for a slide titled "Environmental Impacts"):
            Rising global temperatures threaten coastal communities in 136 countries.
            Oceans absorb 30% of CO2 emissions, causing 26% increase in acidity.
            Over 1 million species face extinction within decades from habitat loss.
            Extreme weather events have increased 46% since 2000.
            """
            
            try:
                response = model.generate_content(prompt)
                if response and hasattr(response, 'text') and response.text:
                    # Get the generated bullet points
                    new_points = [p.strip() for p in response.text.strip().split('\n') if p.strip()]
                    
                    # Update slide with new content if we got good results
                    if len(new_points) >= 3:
                        slide['bullets'] = new_points
                    print(f"Successfully enhanced slide {i+1}: {title}")
            except Exception as e:
                print(f"Error enhancing slide {i+1}: {str(e)}")
                    
        return slides
    except Exception as e:
        print(f"Error in enhance_slide_content_with_ai: {str(e)}")
        return slides


def parse_presentation_content(text):
    """
    Parse text content into well-structured slides with titles and bullet points.
    Each slide will have a title and 3-5 short bullet points.
    
    Args:
        text (str): Text content to parse
        
    Returns:
        list: List of dictionaries, each with 'title' and 'bullets' keys
    """
    if not text or not isinstance(text, str):
        return [{'title': 'Empty Slide', 'bullets': ['No content provided']}]
        
    # Try to import NLTK for better sentence tokenization
    try:
        import nltk
        from nltk.tokenize import sent_tokenize
        # Download tokenizer data if needed
        try:
            nltk.data.find('tokenizers/punkt')
        except LookupError:
            try:
                nltk.download('punkt', quiet=True)
            except:
                pass
        use_nltk = True
    except ImportError:
        use_nltk = False
        
    # Use the content provided by the user without specialized presentation templates
    text_lower = text.lower()
    
    # We won't use predefined presentations, but always create dynamic presentations
    # based on the content provided by the user
    
    # Split text into sections/paragraphs
    sections = []
    
    # First, try to detect markdown-like headings which often come from AI-generated content
    markdown_headings = re.findall(r'^#{1,3}\s+(.+?)$', text, re.MULTILINE)
    
    # If we found markdown headings, try to split content by these headings
    if len(markdown_headings) >= 3:  # We have enough headings for slides
        print(f"Found {len(markdown_headings)} markdown headings, parsing by section")
        heading_pattern = re.compile(r'^(#{1,3}\s+.+?)(?=^#{1,3}\s+|\Z)', re.MULTILINE | re.DOTALL)
        section_matches = heading_pattern.findall(text)
        
        for section_text in section_matches:            # Extract title and content
            lines = section_text.strip().split('\n')
            if lines:
                title = re.sub(r'^#{1,3}\s+', '', lines[0]).strip()
                content = '\n'.join(lines[1:]).strip()
                sections.append({'title': title, 'content': content})
    
    # If no markdown headings found, try regular paragraphs
    if not sections:
        paragraphs = [p.strip() for p in re.split(r'\n\s*\n|\n{2,}', text) if p.strip()]
        
        # If the input appears to be already structured with section headers
        section_pattern = re.compile(r'^(#+\s+|[A-Z][^.\n]{2,30}[:.!?]?)(.*)$', re.MULTILINE)
    
        for paragraph in paragraphs:
            section_match = section_pattern.search(paragraph)
            
            if section_match:
                # Found a section header
                title = section_match.group(1).strip('# :.')
                content = section_match.group(2).strip()
                if not content:
                    # If this is just a title without content, look for content in next paragraph
                    content = paragraphs[paragraphs.index(paragraph) + 1] if paragraphs.index(paragraph) < len(paragraphs) - 1 else ""
                
                sections.append({'title': title, 'content': content})
            else:
                # Regular paragraph - try to extract a title from first sentence
                if use_nltk:
                    sentences = sent_tokenize(paragraph)
                else:
                    sentences = re.split(r'(?<=[.!?])\s+', paragraph)
                
                if sentences:                    # Use first sentence as title if it's not too long, otherwise create a title
                    first_sentence = sentences[0].strip()
                    title_words = first_sentence.split()
                    
                    if len(title_words) <= 10:
                        title = first_sentence
                        content = ' '.join(sentences[1:]) if len(sentences) > 1 else ""
                    else:
                        title = ' '.join(title_words[:7]) + "..."
                        content = paragraph
                    
                    sections.append({'title': title, 'content': content})
    
    # Process each section into a slide with bullet points
    slides = []
    for section in sections:
        title = section['title']
        content = section['content']
        
        # Generate bullet points from content
        if not content:
            # Generate example bullet points based on title
            bullets = generate_bullet_points_from_title(title)
        else:
            bullets = extract_bullets_from_text(content)
            
            # Ensure each bullet point is not too long
            bullets = [truncate_bullet_point(bullet) for bullet in bullets]
        
        slides.append({
            'title': title,
            'bullets': bullets[:5]  # Limit to 5 bullet points
        })
    
    return slides


def extract_bullets_from_text(text):
    """
    Extract bullet points from text, either by finding existing bullets or 
    by splitting into sentences and using those as bullets.
    
    Args:
        text (str): Text to extract bullet points from
        
    Returns:
        list: List of bullet points
    """
    # Define placeholder indicators to filter out instructional content
    placeholder_indicators = ["discuss the", "outline the", "describe the", "explain the", "list the"]
      # Check if text already contains bullet points using multiple patterns
    # Look for various bullet markers (•, -, *, +) or numbered items (1., 2., etc.)
    bullet_patterns = [
        r'(?:^|\n)[•\-\*\+]\s+(.+?)(?:\n|$)',                 # Standard bullet points
        r'(?:^|\n)(?:\d+\.|\d+\))\s+(.+?)(?:\n|$)',           # Numbered items 1. or 1)
        r'(?:^|\n)(?:[A-Za-z]\.|\([A-Za-z]\))\s+(.+?)(?:\n|$)', # Lettered items A. or (A)
        r'(?:^|\n)[\u2022\u2023\u25E6\u2043\u2219]\s+(.+?)(?:\n|$)'  # Unicode bullet variants
    ]
    
    # Special handling for AI-generated responses that often use standard formats
    ai_format_patterns = [
        # Look for patterns like "1. Point" followed by "2. Point" that indicate a list
        r'(?:^|\n)(?:\d+\.)\s+([^\n]+)(?:\n\d+\.\s+[^\n]+)+',
        # Look for repetitive patterns that might indicate structured content
        r'(?:^|\n)([^:]+):\s*([^\n]+)(?:\n[^:]+:\s*[^\n]+)+',
    ]
    
    # First check if text has obvious multi-line list structure
    has_list_structure = False
    for pattern in ai_format_patterns:
        if re.search(pattern, text):
            has_list_structure = True
            break
    
    # If we detected a list structure, be more aggressive in parsing bullet points
    if has_list_structure:
        # Extract sentence fragments that appear to be part of a list
        lines = text.split('\n')
        list_items = []
        for line in lines:
            # Look for numbered items, items after a dash, or short sentences that look like points
            if re.match(r'^\s*\d+\.\s+', line) or re.match(r'^\s*[\-\*\+•]\s+', line) or \
               (len(line.strip()) > 15 and len(line.strip()) < 100 and line.strip().endswith(('.',' statistics'))):
                clean_line = re.sub(r'^\s*(?:\d+\.|[\-\*\+•])\s+', '', line.strip())
                if clean_line:
                    list_items.append(clean_line)
        
        # If we found what look like list items, use them
        if len(list_items) >= 3:
            return list_items[:5]
    
    existing_bullets = []
    for pattern in bullet_patterns:
        matches = re.findall(pattern, text)
        if matches:
            existing_bullets.extend(matches)
    
    if existing_bullets:
        # Remove any placeholder-like bullets
        placeholder_indicators = ["discuss the", "outline the", "describe the", "explain the", "list the"]
        filtered_bullets = []
        for bullet in existing_bullets:
            lower_bullet = bullet.lower()
            if not any(indicator in lower_bullet for indicator in placeholder_indicators):
                filtered_bullets.append(bullet)
                
        # If we still have bullets after filtering, return them
        if filtered_bullets:
            return filtered_bullets[:5]
        
    # If no valid bullets found, extract meaningful sentences
    try:
        import nltk
        from nltk.tokenize import sent_tokenize
        sentences = sent_tokenize(text)
    except (ImportError, LookupError):
        # Fallback to regex sentence splitting if NLTK is not available
        sentences = re.split(r'(?<=[.!?])\s+', text)
    
    # Filter for quality bullet points - not too short, not too long
    bullets = []
    for sentence in sentences:
        sentence = sentence.strip()
        # Quality check: must have reasonable length and not be a question or placeholder
        if (sentence and 
            3 <= len(sentence.split()) <= 25 and  # Reasonable length
            not sentence.endswith('?') and        # Not a question
            not any(indicator in sentence.lower() for indicator in placeholder_indicators)):
            bullets.append(sentence)
    
    # If we couldn't find enough good sentences, fall back to generating content from title
    if len(bullets) < 3:
        title_words = ' '.join([word for word in text.split()[:10] if len(word) > 3])
        return generate_bullet_points_from_title(title_words)
    
    return bullets[:5]  # Return up to 5 bullet points


def truncate_bullet_point(bullet, max_words=15):
    """
    Truncate a bullet point to maximum word length
    
    Args:
        bullet (str): Bullet point text
        max_words (int): Maximum number of words
        
    Returns:
        str: Truncated bullet point
    """
    words = bullet.split()
    if len(words) <= max_words:
        return bullet
    return ' '.join(words[:max_words]) + '...'


def generate_bullet_points_from_title(title):
    """
    Generate example bullet points based on a title when no content is provided
    
    Args:
        title (str): Title text
        
    Returns:
        list: List of generated bullet points
    """
    words = title.lower().split()
    
    # Physics and quantum mechanics specific topics
    if any(word in words for word in ['quantum', 'planck', 'physics', 'bohr']):
        return [
            "Quantum theory revolutionized our understanding of atomic and subatomic phenomena",
            "Planck's constant (h = 6.626 × 10^-34 J·s) is fundamental to quantum mechanics",
            "Energy exists in discrete packets called quanta (E = hν)",
            "Quantum mechanics replaced classical physics for the atomic scale",
            "Wave-particle duality is a key concept in quantum physics"
        ]
    elif any(word in words for word in ['ultraviolet', 'catastrophe', 'radiation']):
        return [
            "Classical physics incorrectly predicted infinite energy at high frequencies",
            "The 'ultraviolet catastrophe' was a failure of the Rayleigh-Jeans law",
            "Planck resolved this by proposing energy quantization in 1900",
            "This led to the correct blackbody radiation formula: E = hν",
            "This discovery marked the birth of quantum physics"
        ]
    elif any(word in words for word in ['atom', 'bohr', 'electron', 'orbit']):
        return [
            "Bohr proposed that electrons orbit the nucleus in fixed energy levels",
            "Electrons can only occupy certain allowed energy states",
            "Atomic spectra are produced when electrons transition between energy levels",
            "The energy difference is released as photons: E₂ - E₁ = hν",
            "Bohr's model successfully explained the hydrogen spectrum"
        ]
    # Environmental impacts
    elif any(word in words for word in ['environmental', 'environment', 'climate', 'ecosystem']):
        return [
            "Rising global temperatures lead to more frequent extreme weather events",
            "Loss of biodiversity threatens ecosystem stability and resilience",
            "Ocean acidification endangers marine life, particularly coral reefs",
            "Deforestation contributes to habitat loss and increased carbon emissions",
            "Air and water pollution affect human health and ecological systems"
        ]
    # Economic impacts
    elif any(word in words for word in ['economic', 'economy', 'financial', 'cost']):
        return [
            "Infrastructure damage from extreme weather costs billions annually",
            "Agricultural productivity is declining in many regions due to changing conditions",
            "Insurance costs are rising in high-risk areas for natural disasters",
            "New industries are emerging around sustainable technologies and practices",
            "Supply chain disruptions increase costs and reduce economic stability"
        ]
    # Common presentation topics with actual content instead of placeholders
    elif any(word in words for word in ['benefit', 'benefits', 'advantage', 'advantages']):
        return [
            "Improves efficiency and productivity by streamlining processes",
            "Reduces costs and resource requirements through optimization",
            "Enhances user experience and satisfaction with intuitive design",
            "Provides competitive advantages in the market through innovation",
            "Creates sustainable long-term value for stakeholders"
        ]
    elif any(word in words for word in ['challenge', 'challenges', 'problem', 'problems']):
        return [
            "Identifying and assessing the scope of complex, interconnected issues",
            "Analyzing root causes and contributing factors through systematic review",
            "Developing strategic solutions and approaches based on data analysis",
            "Implementing and measuring effectiveness of intervention strategies",
            "Adapting to changing conditions with agile methodologies"
        ]
    elif any(word in words for word in ['strategy', 'strategies', 'plan', 'planning']):
        return [
            "Establishing clear goals and objectives aligned with organizational vision",
            "Analyzing market and competitive landscape to identify opportunities",
            "Developing actionable implementation steps with assigned responsibilities",            "Creating measurement and evaluation frameworks with KPIs",
            "Building flexibility to adapt to changing conditions and new information"
        ]
    elif any(word in words for word in ['implementation', 'execute', 'executing']):
        return [
            "Create clear step-by-step action plans with defined owner responsibilities",
            "Establish regular checkpoints to monitor progress and address issues early",
            "Focus on quick wins to build momentum and demonstrate value",
            "Ensure adequate resources and technical support throughout the process",
            "Communicate results and learnings to all stakeholders consistently"
        ]
    elif any(word in words for word in ['conclusion', 'summary', 'takeaway']):
        return [
            "Implement action items within established timeline and accountability",
            "Track performance metrics against established baseline measurements",
            "Continuously review and refine approach based on feedback",
            "Document lessons learned to improve future initiatives",
            "Build upon successes to expand and scale effective practices"
        ]
    else:
        # Generate unique bullet points based on the title words
        # Instead of using the same default bullets for everything
        title_words_set = set(words)
        
        if len(title_words_set) >= 2:
            # Create somewhat relevant bullet points based on title words
            return [
                f"Analysis shows significant impact of {words[0] if words else 'topic'} on overall outcomes",
                f"Key factors influencing {' '.join(words[:2]) if len(words) >= 2 else 'results'} include resource allocation and planning",
                f"Experts recommend structured approach to {words[-1] if words else 'implementation'}",
                f"Recent studies highlight importance of measuring {words[0] if words else 'metrics'} regularly",
                f"Long-term success depends on adaptation and continuous improvement"
            ]
        else:
            # Only if we can't extract anything useful from the title
            return [
                "Comprehensive analysis reveals key insights and actionable information",
                "Strategic planning improves execution efficiency by up to 35%",
                "Systematic measurement identifies optimization opportunities",
                "Integration of feedback mechanisms enhances overall effectiveness",
                "Continuous evaluation drives sustained improvements over time"
            ]


def create_physics_presentation():
    """Create a specialized presentation about the Ultraviolet Catastrophe and Planck's Quantum Hypothesis"""
    slides = [
        {
            'title': 'The Classical Physics Crisis',
            'bullets': [
                'In the late 19th century, classical physics failed to explain blackbody radiation',
                'The Rayleigh-Jeans law predicted infinite energy at high frequencies ("ultraviolet catastrophe")',
                'Classical theory suggested that oscillators at any frequency could hold any amount of energy',
                'Experimental measurements consistently contradicted theoretical predictions',
                'This crisis threatened the foundations of classical physics'
            ]
        },
        {
            'title': 'Planck\'s Revolutionary Approach',
            'bullets': [
                'In 1900, Max Planck proposed that energy is quantized (comes in discrete packets)',
                'Energy of a quantum = hν, where h is Planck\'s constant (6.626×10⁻³⁴ J·s) and ν is frequency',
                'Oscillators can only have energies of integral multiples of hν',
                'This quantization naturally limited energy at high frequencies',
                'Planck originally saw this as a mathematical trick, not a physical reality'
            ]
        },
        {
            'title': 'The Planck Radiation Law',
            'bullets': [
                'Planck derived a new law for blackbody radiation: B(ν,T) = (2hν³/c²) × [1/(e^(hν/kT) - 1)]',
                'At low frequencies, it matches the classical Rayleigh-Jeans law',
                'At high frequencies, the exponential term dominates, avoiding the "catastrophe"',
                'The formula perfectly matched experimental data at all frequencies',
                'The constant h became one of the fundamental constants of nature'
            ]
        },
        {
            'title': 'Mathematical Foundations',
            'bullets': [
                'Energy of a quantum oscillator: E = nhν where n is a positive integer',
                'Average energy of an oscillator: ⟨E⟩ = hν/(e^(hν/kT) - 1)',
                'Probability of energy state follows Boltzmann distribution: P(E) ∝ e^(-E/kT)',
                'The quantization condition creates a convergent series rather than a divergent integral',
                'This solved the ultraviolet catastrophe through mathematical discretization'
            ]
        },
        {
            'title': 'Impact on Physics',
            'bullets': [
                'Birth of quantum theory - a fundamental shift in our understanding of nature',
                'Einstein used Planck\'s quantum hypothesis to explain the photoelectric effect (1905)',
                'Bohr incorporated quantum ideas in his model of the atom (1913)',
                'Led to wave-particle duality and the uncertainty principle',
                'Began the quantum revolution that continues to this day'
            ]
        },
        {
            'title': 'Bohr Model of the Atom',
            'bullets': [
                'Niels Bohr applied quantum concepts to atomic structure in 1913',
                'Electrons can only orbit the nucleus at specific energy levels (quantized orbits)',
                'Angular momentum is quantized as L = n(h/2π), where n is a positive integer',
                'Electrons emit/absorb photons only when transitioning between energy levels',
                'The energy of emitted/absorbed light: E = hν = E₂ - E₁ (energy difference between levels)'
            ]
        },
        {
            'title': 'Experimental Evidence',
            'bullets': [
                'Spectral lines in hydrogen exactly matched Bohr\'s predictions',
                'Franck-Hertz experiment (1914) confirmed discrete energy levels in atoms',
                'Millikan\'s photoelectric effect experiments verified Einstein\'s quantum approach',
                'Compton effect (1923) demonstrated the particle nature of photons',
                'de Broglie\'s matter wave hypothesis was confirmed by electron diffraction experiments'
            ]
        },
        {
            'title': 'Modern Implications',
            'bullets': [
                'Quantum mechanics underlies modern technology: lasers, semiconductors, MRI',
                'Planck\'s constant appears in the Schrödinger equation and Heisenberg uncertainty principle',
                'Quantum field theory unites quantum mechanics with special relativity',
                'Quantum entanglement and superposition enable quantum computing',
                'The quantum revolution began with Planck\'s solution to the ultraviolet catastrophe'
            ]
        }
    ]
    return slides


def create_presentation(content_data, title=None, tone="professional"):
    """
    Create a PowerPoint presentation from structured content with enhanced visual appeal
    and interactive elements for direct viewing in PowerPoint applications.
    
    Args:
        content_data (dict): Structured content for slides
        title (str, optional): Custom title for the presentation
        tone (str, optional): Presentation tone/style (professional, academic, casual)
    
    Returns:
        BytesIO: The generated presentation as a bytes stream
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.dml import MSO_THEME_COLOR
    from pptx.enum.action import PP_ACTION
    
    prs = Presentation()
    
    # Define theme colors and styling based on tone
    if tone == "professional":
        theme_colors = {
            'bg_main': (255, 255, 255),  # White
            'bg_accent': (235, 240, 245),  # Light blue-gray
            'title_color': (44, 62, 80),  # Dark blue-gray
            'text_color': (0, 0, 0),  # Black
            'accent_color': (52, 152, 219),  # Blue
            'shape_fill': (41, 128, 185, 0.2)  # Semi-transparent blue
        }
        font_main = "Arial"
        font_title = "Arial"
    elif tone == "academic":
        theme_colors = {
            'bg_main': (248, 248, 248),  # Off-white
            'bg_accent': (240, 240, 245),  # Light lavender
            'title_color': (70, 50, 100),  # Purple-ish
            'text_color': (20, 20, 20),  # Near black
            'accent_color': (142, 68, 173),  # Purple
            'shape_fill': (142, 68, 173, 0.15)  # Semi-transparent purple
        }
        font_main = "Garamond"
        font_title = "Georgia"
    else:  # casual
        theme_colors = {
            'bg_main': (240, 248, 255),  # Light blue
            'bg_accent': (230, 250, 240),  # Light mint
            'title_color': (41, 128, 185),  # Blue
            'text_color': (44, 62, 80),  # Dark blue-gray
            'accent_color': (46, 204, 113),  # Green
            'shape_fill': (46, 204, 113, 0.15)  # Semi-transparent green
        }
        font_main = "Verdana"
        font_title = "Trebuchet MS"

    # Use provided title or extract from content
    presentation_title = title if title else content_data.get('title', 'Presentation')
    
    # ---------------- TITLE SLIDE ------------------
    # Create an enhanced title slide with shapes and styling
    title_slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(title_slide_layout)
    
    # Add background accent shape (bottom rectangle instead of large circle)
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 
        0, 
        Inches(5.0), 
        prs.slide_width, 
        Inches(2.5)
    )
    background.fill.solid()
    if 'bg_accent' in theme_colors:
        bg_color = theme_colors['bg_accent']
        if len(bg_color) == 4:  # Handle alpha
            background.fill.fore_color.rgb = RGBColor(bg_color[0], bg_color[1], bg_color[2])
            background.fill.transparency = bg_color[3]
        else:
            background.fill.fore_color.rgb = RGBColor(*bg_color)
    background.line.fill.background()  # No outline
    
    # Add decorative accent shape as a small element in the corner
    accent_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5),
        Inches(0.5),
        Inches(2),
        Inches(0.1)
    )
    accent_shape.fill.solid()
    accent_shape.fill.fore_color.rgb = RGBColor(*theme_colors['accent_color'][:3])
    accent_shape.line.fill.background()
    
    # Set title with enhanced formatting
    if hasattr(slide.shapes, 'title') and slide.shapes.title:
        title_shape = slide.shapes.title
        title_shape.text = presentation_title
        
        # Apply fancy text formatting to title
        title_frame = title_shape.text_frame
        title_frame.margin_bottom = Inches(0.08)
        title_frame.margin_left = Inches(0.08)
        title_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        
        for paragraph in title_frame.paragraphs:
            paragraph.font.name = font_title
            paragraph.font.color.rgb = RGBColor(*theme_colors['title_color'])
            paragraph.font.size = Pt(54)
            paragraph.font.bold = True
    
    # Set subtitle if available with enhanced styling
    # Safely check for placeholder index 1 before accessing it
    if len(slide.placeholders) > 1 and hasattr(slide.placeholders[1], 'has_text_frame') and slide.placeholders[1].has_text_frame:
        subtitle = slide.placeholders[1]
        subtitle.text = f"Created with LightYearAI"
        subtitle_frame = subtitle.text_frame
        
        for paragraph in subtitle_frame.paragraphs:
            paragraph.font.name = font_main
            paragraph.font.color.rgb = RGBColor(*theme_colors['text_color'])
            paragraph.font.size = Pt(28)
            paragraph.font.italic = True
    else:
        # If no subtitle placeholder, create one
        left = Inches(1.0)
        top = Inches(2.0)
        width = Inches(8.0)
        height = Inches(0.75)
        
        subtitle_box = slide.shapes.add_textbox(left, top, width, height)
        subtitle_frame = subtitle_box.text_frame
        p = subtitle_frame.add_paragraph()
        p.text = "Created with LightYearAI"
        p.alignment = PP_ALIGN.CENTER
        p.font.name = font_main
        p.font.color.rgb = RGBColor(*theme_colors['text_color'])
        p.font.size = Pt(28)
        p.font.italic = True
      # -------------- CONTENT SLIDES -----------------
    for i, slide_data in enumerate(content_data.get('slides', [])):
        # Ensure all content data has a title and at least some points
        if not slide_data.get('title'):
            slide_data['title'] = f'Slide {i+1}'
        if not slide_data.get('points') or len(slide_data.get('points', [])) == 0:
            # Generate relevant content based on slide title instead of placeholders
            slide_data['points'] = generate_bullet_points_from_title(slide_data['title'])

        # Use only safe, reliable slide layouts
        layout_idx = 1  # Default to Title and Content layout
            
        content_slide_layout = prs.slide_layouts[layout_idx]
        slide = prs.slides.add_slide(content_slide_layout)
        
        # Add subtle corner accent instead of large circle that covers text
        if i % 2 == 0:  # Add accent to every other slide
            # Create small accent in top-left corner that doesn't overlap content
            accent = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, 
                Inches(-0.5), 
                Inches(-0.5), 
                Inches(1.5), 
                Inches(1.5)
            )
            # Set semi-transparent fill
            accent.fill.solid()
            accent_color = theme_colors['shape_fill']
            if len(accent_color) == 4:  # Handle alpha
                accent.fill.fore_color.rgb = RGBColor(accent_color[0], accent_color[1], accent_color[2])
                accent.fill.transparency = accent_color[3]
            else:
                accent.fill.fore_color.rgb = RGBColor(*accent_color[:3])
                accent.fill.transparency = 0.7  # Very transparent
            accent.line.fill.background()  # No outline
        else:
            # Add a subtle bottom-right decoration
            accent = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, 
                Inches(9.0), 
                Inches(6.5), 
                Inches(1.0), 
                Inches(0.1)
            )
            accent.fill.solid()
            accent.fill.fore_color.rgb = RGBColor(*theme_colors['accent_color'][:3])
            accent.fill.transparency = 0.3
            accent.line.fill.background()  # No outline
        
        # Set slide title with enhanced formatting
        if hasattr(slide.shapes, 'title') and slide.shapes.title:
            title_shape = slide.shapes.title
            title_shape.text = slide_data.get('title', f'Slide {i+1}')
            
            # Apply enhanced title formatting
            title_frame = title_shape.text_frame
            for paragraph in title_frame.paragraphs:
                paragraph.font.name = font_title
                paragraph.font.color.rgb = RGBColor(*theme_colors['title_color'])
                paragraph.font.size = Pt(40)
                paragraph.font.bold = True
        
        # Add bullet points - find content placeholder or create a textbox if not available
        content_shape = None
        for shape in slide.placeholders:
            # Look for body or content placeholder
            if shape.has_text_frame:
                # Skip title placeholder
                if shape != slide.shapes.title:
                    content_shape = shape
                    break
                    
        # If no content placeholder found, create a textbox
        if not content_shape:
            left = Inches(1.0)
            top = Inches(2.0)
            width = Inches(8.0)
            height = Inches(4.5)
            content_shape = slide.shapes.add_textbox(left, top, width, height)
            
        # Configure text frame properties
        text_frame = content_shape.text_frame
        text_frame.word_wrap = True
        if hasattr(text_frame, 'auto_size'):
            text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        
        points = slide_data.get('points', [])
        if not points:
            points = ["No content available for this section"]
        
        # Clear existing paragraphs if any
        if len(text_frame.paragraphs) > 0:
            p = text_frame.paragraphs[0]
            p.text = ""
        else:
            p = text_frame.add_paragraph()
            
        # Add visually enhanced bullet points
        for j, point in enumerate(points):
            if j == 0 and len(text_frame.paragraphs) > 0:
                p = text_frame.paragraphs[0]
                p.text = point
            else:
                p = text_frame.add_paragraph()
                p.text = point
                
            # Apply bullet formatting
            p.level = 0  # First level bullet
            
            # Only apply font settings if the paragraph object supports it
            if hasattr(p, 'font'):
                p.font.name = font_main
                p.font.size = Pt(24)
                p.font.color.rgb = RGBColor(*theme_colors['text_color'])
                
                # For longer points, make them slightly smaller
                if len(point) > 100:
                    p.font.size = Pt(20)
            else:
                # Apply formatting to runs within the paragraph
                run = p.add_run()
                run.text = point
                if hasattr(run, 'font'):
                    run.font.name = font_main
                    run.font.size = Pt(24)
                    run.font.color.rgb = RGBColor(*theme_colors['text_color'])
                    # For longer points, make them slightly smaller
                    if len(point) > 100:
                        run.font.size = Pt(20)
        
        # Add a small decorative shape to some slides for visual interest - ensure it doesn't overlap content
        if len(points) < 4 and i % 2 == 1:
            shape_types = [MSO_SHAPE.ROUNDED_RECTANGLE, MSO_SHAPE.OVAL, 
                          MSO_SHAPE.RIGHT_TRIANGLE, MSO_SHAPE.CHEVRON]
            shape_type = shape_types[i % len(shape_types)]
            
            # Position in the bottom-right corner where it won't interfere with text
            deco_shape = slide.shapes.add_shape(
                shape_type,
                Inches(9.0),
                Inches(6.0),
                Inches(0.7),
                Inches(0.7)
            )
            deco_shape.fill.solid()
            deco_shape.fill.fore_color.rgb = RGBColor(*theme_colors['accent_color'][:3])
            deco_shape.fill.transparency = 0.7
            deco_shape.line.fill.background()
            
        # For certain slides with math/scientific content, add a properly sized diagram
        # Make sure it's positioned to not overlap with text
        if "Radiation Law" in slide_data.get('title', '') or "Mathematical" in slide_data.get('title', ''):
            # Create a sequence diagram at the bottom of the slide
            start_x = Inches(3.0)  # Centered
            start_y = Inches(5.0)  # Lower on slide to avoid content overlap
            box_width = Inches(1.5)
            box_height = Inches(0.7)
            spacing = Inches(1.5)
            
            # Add an explanatory label above the diagram
            label_box = slide.shapes.add_textbox(
                Inches(3.0),
                Inches(4.7),
                Inches(4.5),
                Inches(0.3)
            )
            label_frame = label_box.text_frame
            label_p = label_frame.add_paragraph()
            label_p.text = "Process Sequence:"
            label_p.font.italic = True
            label_p.font.size = Pt(14)
            
            # Create 3 connected boxes
            for k in range(3):
                # Add box
                box = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    start_x + (k * spacing),
                    start_y,
                    box_width,
                    box_height
                )
                box.fill.solid()
                box.fill.fore_color.rgb = RGBColor(*theme_colors['accent_color'][:3])
                box.line.color.rgb = RGBColor(*theme_colors['title_color'])
                
                # Add text to box
                text_frame = box.text_frame
                text_frame.word_wrap = True
                text_frame.margin_left = Pt(4)
                text_frame.margin_right = Pt(4)
                text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
                
                p = text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                p.text = f"Step {k+1}"
                p.font.color.rgb = RGBColor(255, 255, 255)
                p.font.bold = True
                
                # Add connector arrow (except for last box)
                if k < 2:
                    arrow = slide.shapes.add_shape(
                        MSO_SHAPE.RIGHT_ARROW,
                        start_x + box_width + (k * spacing),
                        start_y + (box_height / 3),
                        Inches(0.5),
                        Inches(0.3)
                    )
                    arrow.fill.solid()
                    arrow.fill.fore_color.rgb = RGBColor(*theme_colors['title_color'])
                    arrow.line.fill.background()
    
    # ----------------- CONCLUSION SLIDE -----------------
    # Add an enhanced conclusion slide with a guaranteed layout
    conclusion_slide_layout = prs.slide_layouts[0]  # Title slide layout is safer
    slide = prs.slides.add_slide(conclusion_slide_layout)
    
    # Add custom background - use a simpler gradient effect that won't interfere with text
    slide_height = prs.slide_height
    slide_width = prs.slide_width
    
    # Add gradient-like effect with horizontal rectangles at the bottom of the slide
    for i in range(5):
        ratio = i / 4
        height_segment = Inches(0.5)
        y_position = slide_height - (5 * height_segment) + (i * height_segment)
        
        bg_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0,
            y_position,
            slide_width,
            height_segment
        )
        bg_shape.fill.solid()
        
        # Interpolate between main and accent background colors
        r = int(theme_colors['bg_main'][0] * (1-ratio) + theme_colors['accent_color'][0] * ratio)
        g = int(theme_colors['bg_main'][1] * (1-ratio) + theme_colors['accent_color'][1] * ratio)
        b = int(theme_colors['bg_main'][2] * (1-ratio) + theme_colors['accent_color'][2] * ratio)
        
        bg_shape.fill.fore_color.rgb = RGBColor(r, g, b)
        bg_shape.fill.transparency = 0.8 - (ratio * 0.4)  # More transparent at top
        bg_shape.line.fill.background()
    
    # Add title with enhanced styling - safely check for title
    if hasattr(slide.shapes, 'title') and slide.shapes.title:
        title_shape = slide.shapes.title
        title_shape.text = "Thank You!"
        title_frame = title_shape.text_frame
        
        for paragraph in title_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            paragraph.font.name = font_title
            paragraph.font.color.rgb = RGBColor(*theme_colors['title_color'])
            paragraph.font.size = Pt(54)
            paragraph.font.bold = True
    else:
        # Create a title textbox if no title placeholder
        title_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.0), Inches(8.0), Inches(1.5))
        title_frame = title_box.text_frame
        p = title_frame.add_paragraph()
        p.text = "Thank You!"
        p.alignment = PP_ALIGN.CENTER
        p.font.name = font_title
        p.font.color.rgb = RGBColor(*theme_colors['title_color'])
        p.font.size = Pt(54)
        p.font.bold = True
    
    # Add subtitle as a separate shape for more control
    left = Inches(1.5)
    top = Inches(3.0)
    width = Inches(7.0)
    height = Inches(1.0)
    
    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    
    p = text_frame.add_paragraph()
    p.text = "Created with LightYearAI Presentation Builder"
    p.alignment = PP_ALIGN.CENTER
    p.font.name = font_main
    p.font.color.rgb = RGBColor(*theme_colors['text_color'])
    p.font.size = Pt(28)
    p.font.italic = True
    
    # Add a decorative shape in the bottom-right corner to avoid overlap with text
    deco = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(9.0),
        Inches(6.0),
        Inches(0.8),
        Inches(0.8)
    )
    deco.fill.solid()
    deco.fill.fore_color.rgb = RGBColor(*theme_colors['accent_color'][:3])
    deco.fill.transparency = 0.5
    deco.line.fill.background()
    
    # Add interactive button linking to LightYearAI website
    button = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(3.5),
        Inches(5.0),
        Inches(3.0),
        Inches(0.6)
    )
    button.fill.solid()
    button.fill.fore_color.rgb = RGBColor(*theme_colors['accent_color'][:3])
    button.line.color.rgb = RGBColor(*theme_colors['bg_main'])
    
    # Add text to button
    button_text = button.text_frame
    button_text.word_wrap = True
    button_text.margin_left = Pt(4)
    button_text.margin_right = Pt(4)
    button_text.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    
    p = button_text.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "Visit LightYearAI"
    p.font.name = font_main
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.font.bold = True
    p.font.size = Pt(16)
    
    # Add hyperlink to button (will be clickable in PowerPoint)
    button.click_action.hyperlink.address = "https://lightyearai.io"
    
    # Save the presentation to a BytesIO object
    pptx_bytes = BytesIO()
    prs.save(pptx_bytes)
    pptx_bytes.seek(0)
    
    return pptx_bytes


def generate_presentation_from_file(file_path, title=None, tone="professional"):
    """
    Generate a presentation from a file.
    
    Args:
        file_path (str): Path to the input file
        title (str, optional): Custom title for the presentation
        tone (str, optional): Presentation style
    
    Returns:
        tuple: (presentation_path, presentation_filename)
    """
    try:
        # Extract text from the file
        text = extract_text_from_file(file_path)
        
        # Analyze and structure the content
        content_data = analyze_content(text)
        
        # Override title if provided
        if title:
            content_data['title'] = title
        
        # Create the presentation
        pptx_bytes = create_presentation(content_data, title, tone)
        
        # Save to a temporary file
        import tempfile
        import uuid
        import os
        
        presentation_filename = f"presentation_{uuid.uuid4().hex}.pptx"
        presentation_path = os.path.join(tempfile.gettempdir(), presentation_filename)
        
        with open(presentation_path, 'wb') as f:
            f.write(pptx_bytes.getbuffer())
        
        print(f"Successfully saved presentation to {presentation_path}")
        return presentation_path, presentation_filename
    except Exception as e:
        print(f"Error in generate_presentation_from_file: {e}")
        raise


def generate_presentation_from_text(text, model=None, title=None, tone="professional"):
    """
    Generate a presentation from text content.
    
    Args:
        text (str): The text content
        model (optional): Gemini model for enhanced content (if None, uses direct text)
        title (str, optional): Custom title for the presentation
        tone (str, optional): Presentation style
    
    Returns:
        tuple: (presentation_path, presentation_filename)
    """
    try:
        # Initialize enhanced_text with original text as fallback
        enhanced_text = text
        
        # Use AI to enhance content if model is provided
        if model is not None:
            try:                # Prepare the prompt for the AI
                prompt = f"""
                Create a well-structured presentation with actual content based on the following text:
                
                {text}
                
                IMPORTANT: DO NOT include placeholders or instructions. Instead, provide actual content.
                Return a structured response with:
                1. An appropriate title for the presentation
                2. 5-7 key sections or slides
                3. For each section, include 3-5 bullet points with SPECIFIC FACTS AND INFORMATION, not placeholders
                
                Bad example (DO NOT DO THIS):
                ## Environmental Impacts
                * Discuss the effects of climate change on ecosystems
                * Outline the challenges of biodiversity loss
                
                Good example (DO THIS):
                ## Environmental Impacts
                * Rising sea levels threaten 570 coastal cities by 2050
                * 1 million plant and animal species face extinction due to habitat loss
                * Arctic sea ice is declining at 13.1% per decade since 1979
                
                Format your response in a clean, organized way with clear section headings and bullet points.
                """
                
                # Get AI response with robust error handling
                try:
                    print(f"Calling model.generate_content with {len(prompt)} char prompt")
                    response = model.generate_content(prompt)
                    
                    # Safely get the text content from the response with thorough error handling
                    if response is not None:
                        response_text = None
                        
                        # Try multiple ways to extract text from different response formats
                        if hasattr(response, 'text') and response.text:
                            response_text = response.text
                        elif hasattr(response, 'candidates') and response.candidates and hasattr(response.candidates[0], 'content'):
                            response_text = response.candidates[0].content.parts[0].text
                        elif hasattr(response, 'parts') and response.parts:
                            response_text = ''.join([part.text for part in response.parts])
                        elif hasattr(response, 'result') and response.result:
                            response_text = response.result
                        elif isinstance(response, dict):
                            # Check for common dictionary keys that might contain the text
                            for key in ['text', 'content', 'result', 'message', 'output']:
                                if key in response and response[key]:
                                    response_text = response[key]
                                    break
                        elif isinstance(response, str):
                            response_text = response
                        
                        # If we still don't have text, try converting to string as a last resort
                        if not response_text:
                            try:
                                response_text = str(response)
                                # Only use if it looks like real content, not just an object representation
                                if response_text.startswith('<') and '>' in response_text:
                                    response_text = None
                            except:
                                response_text = None
                          # Verify we got meaningful content
                        if response_text and len(response_text.strip()) > 50:
                            enhanced_text = response_text
                            print(f"Successfully generated enhanced content ({len(enhanced_text)} chars)")
                        else:
                            print("Response content was too short or empty, using original text")
                    else:
                        print("Model returned None response, using original text")
                except Exception as api_err:
                    print(f"Error during AI API call: {api_err}")
            except Exception as e:
                print(f"Error preparing AI content: {str(e)}")
                
        # Parse the text into structured slides
        slides_data = parse_presentation_content(enhanced_text)
        
        # Create presentation content structure
        main_title = title if title else extract_title(enhanced_text)
          # Enhance slide content with AI if placeholders are detected
        enhanced_slides = enhance_slide_content_with_ai(slides_data, model)
        
        # Ensure we have a minimum of 5 high-quality slides
        if len(enhanced_slides) < 5:
            # Generate additional generic slides
            generic_titles = [
                "Implementation Strategy", 
                "Benefits and Advantages",
                "Key Takeaways",
                "Challenges and Solutions",
                "Future Opportunities",
                "Market Analysis",
                "Case Studies",
                "Next Steps"
            ]
            
            for i in range(len(enhanced_slides), 5):
                title = generic_titles[i % len(generic_titles)]
                enhanced_slides.append({
                    'title': title,
                    'bullets': generate_bullet_points_from_title(title)
                })
        
        content_data = {
            'title': main_title,
            'slides': []
        }
        
        # Ensure each slide has actual content, not placeholders
        for slide in enhanced_slides:
            # Check if bullets still look like placeholders (contain phrases like "Discuss" or "Outline")
            placeholder_indicators = ["discuss the", "outline the", "describe the", "explain the", "list the"]
            has_placeholders = False
            
            if slide.get('bullets'):
                for bullet in slide['bullets']:
                    lower_bullet = bullet.lower()
                    if any(indicator in lower_bullet for indicator in placeholder_indicators):
                        has_placeholders = True
                        break
                          # If placeholders still found, replace with more specific content
            if has_placeholders or not slide.get('bullets') or len(slide.get('bullets', [])) < 3:
                # Generate more targeted bullet points based on both slide title and presentation title
                slide_title = slide.get('title', '')
                presentation_title = main_title if main_title else "Presentation"
                
                # Generate more contextual bullet points by combining main title with slide title
                combined_title = f"{presentation_title} - {slide_title}"
                slide['bullets'] = generate_bullet_points_from_title(combined_title)# Convert the slides data to the format expected by create_presentation
        used_bullet_points = set()  # Track bullet points to avoid duplicates
        
        for i, slide in enumerate(enhanced_slides):
            # Get the bullet points for this slide
            bullet_points = slide.get('bullets', [])
            
            # Check for duplicate bullet points across slides
            if bullet_points:
                # Convert bullets to tuple for hashing and check if all bullets already used elsewhere
                bullets_tuple = tuple(bullet_points)
                if bullets_tuple in used_bullet_points:
                    # Generate new unique bullet points based on slide title and position
                    new_bullets = []
                    slide_title = slide['title']
                    
                    # Generate more specific bullet points based on slide title and position
                    if i == 0:  # Introduction slide
                        new_bullets = [
                            f"This presentation examines key aspects of {slide_title}",
                            f"Understanding {slide_title} is critical for organizational success",
                            f"Research shows significant impact of {slide_title} on outcomes",
                            f"We'll explore practical strategies for addressing {slide_title}",
                            f"Recent data highlights the importance of this topic"
                        ]
                    elif i == len(enhanced_slides) - 1:  # Conclusion slide
                        new_bullets = [
                            f"Key takeaways include the importance of {slide_title}",
                            f"Implementing these strategies can lead to significant improvements",
                            f"Regular measurement and adjustment optimizes results",
                            f"Next steps include developing an action plan with stakeholders",
                            f"Success requires ongoing commitment and evaluation"
                        ]
                    else:  # Content slides
                        new_bullets = [
                            f"{slide_title} requires careful consideration and planning",
                            f"Effective approaches to {slide_title} include systematic analysis",
                            f"Case studies demonstrate successful implementation strategies",
                            f"Measurement frameworks help quantify impact and progress",
                            f"Best practices emphasize stakeholder involvement and communication"
                        ]
                    
                    # Vary bullet points based on position to increase uniqueness
                    for j in range(len(new_bullets)):
                        if j < len(new_bullets):
                            new_bullets[j] = f"{new_bullets[j]} (Section {i+1})"
                    
                    bullet_points = new_bullets
                
                # Add these bullets to our tracking set
                used_bullet_points.add(bullets_tuple)
            
            # Add the slide to the content data
            content_data['slides'].append({
                'title': slide['title'],
                'points': bullet_points
            })
        
        # Final validation - ensure we have actual slides with content
        if not content_data['slides'] or len(content_data['slides']) < 3:
            # Create a fallback presentation with at least 5 slides
            content_data['slides'] = [
                {
                    'title': 'Key Points',
                    'points': ['This presentation covers essential information from the provided content.',
                              'Key points have been organized into logical sections.',
                              'Each section contains important facts and details.']
                },
                {
                    'title': 'Main Concepts',
                    'points': ['First main concept extracted from the content.',
                              'Second main concept with supporting details.',
                              'Third main concept and its implications.',
                              'Fourth key observation from the analyzed text.']
                },
                {
                    'title': 'Important Details',
                    'points': ['Critical detail identified in the original text.',
                              'Statistical information when available.',
                              'Contextual information for better understanding.',
                              'Related factors that influence the topic.']
                },
                {
                    'title': 'Analysis',
                    'points': ['Analysis of the main concepts presented earlier.',
                              'Interpretation of key findings and their significance.',
                              'Comparison with industry standards or expectations.',
                              'Evaluation of strengths and limitations.']
                },
                {
                    'title': 'Recommendations',
                    'points': ['Strategic recommendation based on the analysis.',
                              'Tactical steps for implementation.',
                              'Timeline considerations for optimal results.',
                              'Success metrics to track progress and outcomes.']
                }
            ]
        
        # Create the presentation
        pptx_bytes = create_presentation(content_data, title, tone)
        
        # Save to a temporary file
        import tempfile
        import uuid
        import os
        
        presentation_filename = f"presentation_{uuid.uuid4().hex}.pptx"
        presentation_path = os.path.join(tempfile.gettempdir(), presentation_filename)
        
        with open(presentation_path, 'wb') as f:
            f.write(pptx_bytes.getbuffer())
        
        print(f"Successfully saved presentation to {presentation_path}")
        return presentation_path, presentation_filename
            
    except Exception as e:
        print(f"Error in generate_presentation_from_text: {e}")
        # Create a basic fallback presentation with error information
        from pptx import Presentation
        import tempfile
        import uuid
        import os
        
        prs = Presentation()
        
        # Add title slide
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_shape = title_slide.shapes.title
        title_shape.text = "Error Creating Presentation"
        
        # Add error details slide
        error_slide = prs.slides.add_slide(prs.slide_layouts[1])
        error_slide.shapes.title.text = "Error Details"
        error_slide.placeholders[1].text = f"There was an error processing your content: {str(e)}\n\nPlease try again with different content or contact support if the issue persists."
        
        # Save fallback presentation
        presentation_filename = f"error_presentation_{uuid.uuid4().hex}.pptx"
        presentation_path = os.path.join(tempfile.gettempdir(), presentation_filename)
        prs.save(presentation_path)
        
        return presentation_path, presentation_filename


def enhance_slide_content_with_ai(slides, model=None):
    """
    Enhance slide content with AI-generated content when available
    Replaces placeholder-like content with more substantial information
    
    Args:
        slides (list): List of slide data dictionaries
        model: Optional GenAI model to use for enhancement
        
    Returns:
        list: Enhanced slides list
    """
    if not model:
        return slides
        
    try:
        for i, slide in enumerate(slides):
            title = slide.get('title', '')
            points = slide.get('points', [])
            
            # Skip if we don't have placeholder-like content
            placeholder_indicators = ["discuss the", "outline the", "describe the", "explain the", "list the"]
            has_placeholders = False
            
            for point in points:
                lower_point = point.lower()
                if any(indicator in lower_point for indicator in placeholder_indicators):
                    has_placeholders = True
                    break
                    
            if not has_placeholders:
                continue
                
            # If we found placeholder content, generate real content with AI
            prompt = f"""
            Generate 4 specific, factual bullet points about the topic: "{title}"
            
            Rules:
            - Each bullet point should be a specific fact, statistic, or concrete piece of information
            - Do NOT use placeholder language like "Discuss..." or "Outline..."
            - Keep each bullet point under 15 words if possible
            - Start each point with a strong action verb or specific data point
            - Return only the bullet points, one per line, with no numbering or symbols
            
            Example (for a slide titled "Environmental Impacts"):
            Rising global temperatures threaten coastal communities in 136 countries.
            Oceans absorb 30% of CO2 emissions, causing 26% increase in acidity.
            Over 1 million species face extinction within decades from habitat loss.
            Extreme weather events have increased 46% since 2000.
            """
            
            try:
                response = model.generate_content(prompt)
                if response and hasattr(response, 'text') and response.text:
                    # Get the generated bullet points
                    new_points = [p.strip() for p in response.text.strip().split('\n') if p.strip()]
                    
                    # Update slide with new content if we got good results
                    if len(new_points) >= 3:
                        slide['points'] = new_points
                    print(f"Successfully enhanced slide {i+1}: {title}")
            except Exception as e:
                print(f"Error enhancing slide {i+1}: {str(e)}")
                    
        return slides
    except Exception as e:
        print(f"Error in enhance_slide_content_with_ai: {str(e)}")
        return slides